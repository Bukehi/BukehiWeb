from pptx import Presentation
from pptx.util import Cm
import cv2
import os


def audio(filepath, name):
    # 创建一个新的PowerPoint文档
    prs = Presentation('./ex/template.pptx')
    save_poster_temp = './ex/audio.png'

    # 设置幻灯片的大小
    # prs.slide_width = Cm(33.867)
    # prs.slide_height = Cm(19.05)

    # 添加一个空白模板
    # slide_layout = prs.slide_layouts[6]   # 选择带有标题和内容的幻灯片布局
    slide = prs.slides[0]

    # 添加一个音频
    try:
        audio_file1 = filepath[0]   # 音频文件路径
        slide.shapes.add_movie(audio_file1, Cm(3.5), Cm(
            2.5), Cm(26.7), Cm(1.9), mime_type='audio/mp3', poster_frame_image=save_poster_temp)
    except:
        pass
    try:
        audio_file2 = filepath[1]
        slide.shapes.add_movie(audio_file2, Cm(3.5), Cm(
            5.44), Cm(26.7), Cm(1.9), mime_type='audio/mp3', poster_frame_image=save_poster_temp)
    except:
        pass
    try:
        audio_file3 = filepath[2]
        slide.shapes.add_movie(audio_file3, Cm(3.5), Cm(
            8.38), Cm(26.7), Cm(1.9), mime_type='audio/mp3', poster_frame_image=save_poster_temp)
    except:
        pass

    # 保存演示文稿
    prs.save(f"{name}.pptx")


def movie(filepath, name):
    # 创建一个新的PowerPoint文档
    prs = Presentation()

    # 设置幻灯片的大小
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)

    # 添加一个空白模板
    slide_layout = prs.slide_layouts[6]   # 选择带有标题和内容的幻灯片布局
    slide = prs.slides.add_slide(slide_layout)

    # 添加一个视频
    try:
        video_path = filepath
        cap = cv2.VideoCapture(video_path)
        size = (int(cap.get(cv2.CAP_PROP_FRAME_WIDTH)),
                int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT)))
        if size[0] <= 960:
            width = 0.052875*size[0]
            hight = 0.052875*size[1]
        elif size[0] > 960:
            width = 0.052875*size[0]/2
            hight = 0.052875*size[1]/2
        cap.set(cv2.CAP_PROP_POS_FRAMES, 0)  # 设置要获取的帧
        ret, frame = cap.read()
        cv2.imwrite('save_poster_temp.jpg', frame)
        cap.release()
        slide.shapes.add_movie(video_path, Cm(0), Cm(0), Cm(width), Cm(
            hight), mime_type='video/mp4', poster_frame_image='save_poster_temp.jpg')
        os.remove('save_poster_temp.jpg')
    except:
        pass
    prs.save(f"{name}.pptx")


if __name__ == "__main__":
    audio(['C:/music/chun/1media1.mp3'], '纯音1')
